#!/usr/bin/env python
from pptx import Presentation

p = Presentation()
slide = p.slides.add_slide(p.slidelayouts[0])
slide.shapes.title.text = "Hello Title"
slide.shapes.placeholders[1].text = "hello subtitle"

p.save('pptx_test.pptx')
